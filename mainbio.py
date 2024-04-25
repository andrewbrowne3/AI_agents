from dotenv import load_dotenv
from llama_index.core import (
    VectorStoreIndex,
    SimpleKeywordTableIndex,
    SimpleDirectoryReader,
)
from llama_index.core import SummaryIndex
from llama_index.core.schema import IndexNode
from llama_index.core.tools import QueryEngineTool, ToolMetadata
from llama_index.core.callbacks import CallbackManager
from llama_index.llms.openai import OpenAI
from llama_parse import LlamaParse
import os
from llama_index.core.memory.chat_memory_buffer import ChatMemoryBuffer
from llama_index.core.tools import FunctionTool
from llama_index.core.agent import (
    FunctionCallingAgentWorker,
    AgentRunner,
    ReActAgent,
)
#from memory import ChatMemoryBuffer


load_dotenv()
os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY")

reviews = [
    "avogadros_number",
    "Point_of_care_review",
    "Biosensor_review",
    "Accreditation",
    "chemdoodlejs"
 
]

# Example tokenizer function
def get_tokenizer():
    # This is a simplistic tokenizer and should be replaced with one appropriate for your application
    return lambda text: text.split()

from pathlib import Path

import requests


# Load all wiki documents
city_docs = {}
for review in reviews:
    city_docs[review] = SimpleDirectoryReader(
        input_files=[f"data/Papers/{review}.pdf"]
    ).load_data()


llm = OpenAI(temperature=0, model="gpt-3.5-turbo")
callback_manager = CallbackManager([])


from llama_index.agent.openai import OpenAIAgent
from llama_index.core import load_index_from_storage, StorageContext
from llama_index.core.node_parser import SentenceSplitter
import os

node_parser = SentenceSplitter()

# Build agents dictionary
query_engine_tools = []

for idx, review in enumerate(reviews):
    nodes = node_parser.get_nodes_from_documents(city_docs[review])

    if not os.path.exists(f"./data/Papers/{review}"):
        # build vector index
        vector_index = VectorStoreIndex(
            nodes, callback_manager=callback_manager
        )
        vector_index.storage_context.persist(
            persist_dir=f"./data/Papers/{review}"
        )
    else:
        vector_index = load_index_from_storage(
            StorageContext.from_defaults(persist_dir=f"./data/Papers/{review}"),
            callback_manager=callback_manager,
        )
    # define query engines
    vector_query_engine = vector_index.as_query_engine(llm=llm)
    print(query_engine_tools)
    # define tools
    query_engine_tools.append(
        QueryEngineTool(
            query_engine=vector_query_engine,
            metadata=ToolMetadata(
                name=f"vector_tool_{review}",
                description=(
                    "Useful for questions related to specific aspects of"
                    f" {review} The background scientific explanation of the conclusion The transducer of the biosenor etc."
                ),
            ),
        )
    )
    print(query_engine_tools)
# this part explains how to run the query 
from llama_index.core.agent import AgentRunner
from llama_index.agent.openai import OpenAIAgentWorker, OpenAIAgent
from llama_index.agent.openai import OpenAIAgentWorker
def test():
    return ""

def multiply(a: int, b: int) -> int:
    """Multiply two integers and returns the result integer"""
    return a * b


def add(a: int, b: int) -> int:
    """Add two integers and returns the result integer"""
    return a + b


def subtract(a: int, b: int) -> int:
    """Subtract two integers and returns the result integer"""
    return a - b

from rdkit import Chem
from rdkit.Chem import Draw

def smiles_to_mol(smiles):
    """
    Convert a SMILES string to an RDKit molecule object.
    
    Args:
    - smiles (str): A SMILES string representing the chemical structure.
    
    Returns:
    - RDKit Mol object or None if invalid SMILES.
    """
    mol = Chem.MolFromSmiles(smiles)
    if mol is None:
        print("Invalid SMILES string provided.")
    return mol

def draw_molecule(mol, file_name):
    """
    Draw a molecule and save the image to a file using RDKit.
    
    Args:
    - mol (RDKit Mol): A molecule object. 
    - file_name (str): The file name to save the image.
    
    Returns:
    - None; saves the image file directly.
    """
    if mol is not None:
        img = Draw.MolToImage(mol)
        img.save(file_name)
        print(f"Image saved as {file_name}")
    else:
        print("No valid molecule provided for drawing.")
import subprocess
import sys

from pptx import Presentation

def create_powerpoint_presentation(title_text="Default Title", content_text="Default Content"):
    """
    Create a PowerPoint presentation with one slide.

    Args:
    title_text (str): Text for the title of the slide.
    content_text (str): Text for the content of the slide.
    """
    # Create a presentation object
    prs = Presentation()
    
    # Add a slide with the title and content layout
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Set the title and content
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = title_text  # Use the function parameter for the title
    content.text = content_text  # Use the function parameter for the content
    
    # Save the presentation
    prs.save('demo.pptx')
def open_google_chrome(url='https://www.google.com'):
    """
    Opens Google Chrome with a specified URL.

    Args:
    - url (str): URL to open in Google Chrome. Defaults to Google's homepage.
    """
    try:
        if sys.platform.startswith('linux') or sys.platform.startswith('linux2'):
            # Linux
            subprocess.run(['google-chrome', url])
        elif sys.platform.startswith('darwin'):
            # macOS
            subprocess.run(['open', '-a', 'Google Chrome', url])
        elif sys.platform.startswith('win32'):
            # Windows
            subprocess.run(['start', 'chrome', url], shell=True)
    except Exception as e:
        print(f"Failed to open Google Chrome: {e}")

import sympy as sp
import matplotlib.pyplot as plt
import numpy as np


def differentiate_and_plot(expression, variable):
    """
    Differentiates a mathematical expression with respect to a given variable and plots it.
    
    Args:
    expression (str): The mathematical expression to differentiate.
    variable (str): The variable with respect to which differentiation is performed.
    
    Returns:
    str: The derivative of the expression.
    """
    x = sp.symbols(variable)
    f = sp.sympify(expression)
    derivative = sp.diff(f, x)

    # Plotting
    sp.plot(derivative, title=f"Derivative of {expression}", legend=True, show=True)

    return str(derivative)

def calculate_limit_and_plot(expression, variable, point):
    """
    Calculates the limit of an expression as the variable approaches a point and plots the function.
    
    Args:
    expression (str): The mathematical expression for which the limit is calculated.
    variable (str): The variable in the expression.
    point (float or str): The point at which the limit is evaluated.
    
    Returns:
    str: The calculated limit.
    """
    x = sp.symbols(variable)
    f = sp.sympify(expression)
    limit_value = sp.limit(f, x, point)

    # Plotting
    p = sp.plot(f, title=f"Function {expression} and its Limit as {variable} -> {point}", legend=True, show=False)
    p.ylim = [float(limit_value) - 1, float(limit_value) + 1]  # Adjust y limits to focus on the limit point
    p.show()

    return str(limit_value)

def integrate_and_plot(expression, variable):
    """
    Integrates a mathematical expression with respect to a given variable and plots it.
    
    Args:
    expression (str): The mathematical expression to integrate.
    variable (str): The variable with respect to which integration is performed.
    
    Returns:
    str: The integral of the expression.
    """
    x = sp.symbols(variable)
    f = sp.sympify(expression)
    integral = sp.integrate(f, x)

    # Plotting
    sp.plot(integral, title=f"Integral of {expression}", legend=True, show=True)

    return str(integral)

import json

def quiz_to_enhanced_json(quiz_questions):
    """
    Converts a list of quiz questions with additional attributes into a JSON format.

    Args:
    quiz_questions (list): A list of dictionaries, each containing detailed quiz question attributes.

    Returns:
    str: A JSON string representing the quiz questions.
    """
    # Construct the JSON object from the quiz questions
    quiz_dict = {
        "quiz": []
    }
    for item in quiz_questions:
        question_format = {
            "text": item["text"],
            "type": item.get("type", "MCQ"),  # Default to "MCQ" if not provided
            "choices": item["choices"],
            "correct_answer": item["correct_answer"],
            "marks": item.get("marks", 1),  # Default to 1 if not provided
            "difficulty": item.get("difficulty", 0.5),  # Default to 0.5 if not provided
            "discrimination": item.get("discrimination", 0.5),  # Default to 0.5 if not provided
            "explanation": item.get("explanation", ""),
            "course": item.get("course", 0),  # Default to 0 if not provided
            "topics": item.get("topics", [])
        }
        quiz_dict["quiz"].append(question_format)

    # Convert the Python dictionary to a JSON string
    json_output = json.dumps(quiz_dict, indent=4)
    return json_output

# Example usage:
# smiles_str = "CCO"
# mol = smiles_to_mol(smiles_str)
# img = draw_molecule(mol)
# img.show()
smiles_string_to_rdkitobj_tool=FunctionTool.from_defaults(fn=smiles_to_mol)
draw_moleculetool=FunctionTool.from_defaults(fn=draw_molecule)
websitetool = FunctionTool.from_defaults(fn=open_google_chrome)
multiply_tool = FunctionTool.from_defaults(fn=multiply)
add_tool = FunctionTool.from_defaults(fn=add)
subtract_tool = FunctionTool.from_defaults(fn=subtract)
pptcreation = FunctionTool.from_defaults(fn=create_powerpoint_presentation)
differentiation = FunctionTool.from_defaults(fn=differentiate_and_plot)
limits = FunctionTool.from_defaults(fn=calculate_limit_and_plot)
integral = FunctionTool.from_defaults(fn=integrate_and_plot)
quiz = FunctionTool.from_defaults(fn=quiz_to_enhanced_json)


y = 0
while y < 10:
    x = int(input("1 or 2"))
    openai_step_engine = OpenAIAgentWorker.from_tools(
        query_engine_tools, llm=llm, verbose=True
    )
    agent_worker = FunctionCallingAgentWorker.from_tools(
        [multiply_tool, add_tool, subtract_tool,draw_moleculetool,
        smiles_string_to_rdkitobj_tool,websitetool,pptcreation,
        differentiation,limits,integral,quiz],
        llm=llm,
        verbose=True,
        allow_parallel_tool_calls=False
    )
    agent = AgentRunner(agent_worker)

    agent = AgentRunner(openai_step_engine) if x==1 else AgentRunner(agent_worker)
    # # alternative
    #def test():
        #print("Works")
    #agent = OpenAIAgent.from_tools(query_engine_tools, llm=llm, verbose=True,function=test)

    response = agent.chat(
        input("What do you want to know about")
    )

    y+=1
